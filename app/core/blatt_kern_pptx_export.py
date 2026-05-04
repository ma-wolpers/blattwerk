"""PPTX export: renders each presentation slide (or worksheet page) as an image
and packages them into a python-pptx presentation."""

from __future__ import annotations

import io
import tempfile
from pathlib import Path

import fitz  # PyMuPDF
from PIL import Image

from .build_requests import WorksheetBuildRequest, build_worksheet_from_request


# Slide dimensions in cm for each page format
_PAGE_SIZE_CM = {
    "presentation_16_9":  (33.867, 19.05),
    "presentation_16_10": (30.48,  19.05),
    "presentation_4_3":   (25.4,   19.05),
    "a4_portrait":        (21.0,   29.7),
    "a5_landscape":       (21.0,   14.8),
}

# 1 cm = 360000 EMU (English Metric Units, the unit python-pptx uses)
_CM_TO_EMU = 360_000


def _cm_to_emu(value_cm: float) -> int:
    return round(value_cm * _CM_TO_EMU)


def build_presentation_pptx(
    input_path: Path,
    output_path: Path,
    page_format: str = "presentation_16_9",
    print_profile: str = "standard",
    design=None,
    include_solutions: bool = False,
    black_screen_mode: str = "none",
    presentation_section_separator: str = "dot",
    presentation_hide_future_sections: bool = False,
    metadata_defaults: dict | None = None,
    copyright_text_override: str | None = None,
    render_dpi: int = 200,
) -> None:
    """Render every page of the document to an image and write a .pptx file.

    Each rendered page becomes one slide.  The slide dimensions match the chosen
    page_format so the images fill the slides exactly without borders.
    """
    try:
        from pptx import Presentation
        from pptx.util import Emu
    except ImportError as exc:
        raise RuntimeError(
            "python-pptx ist nicht installiert. "
            "Installiere es mit: pip install python-pptx"
        ) from exc

    # 1. Render to a temporary PDF via the normal pipeline
    with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
        temp_pdf_path = Path(tmp.name)

    try:
        build_worksheet_from_request(
            WorksheetBuildRequest(
                input_path=input_path,
                output_path=temp_pdf_path,
                include_solutions=include_solutions,
                page_format=page_format,
                print_profile=print_profile,
                design=design,
                metadata_defaults=metadata_defaults or {},
                copyright_text_override=copyright_text_override,
                black_screen_mode=black_screen_mode,
                presentation_section_separator=presentation_section_separator,
                presentation_hide_future_sections=presentation_hide_future_sections,
            )
        )

        # 2. Convert each PDF page to a PIL image
        slide_images: list[Image.Image] = []
        with fitz.open(temp_pdf_path) as doc:
            for page_index in range(len(doc)):
                page = doc.load_page(page_index)
                pix = page.get_pixmap(dpi=render_dpi, alpha=False)
                img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
                slide_images.append(img)

    finally:
        try:
            temp_pdf_path.unlink(missing_ok=True)
        except Exception:
            pass

    if not slide_images:
        raise ValueError("Das Dokument hat keine Seiten – PPTX kann nicht erstellt werden.")

    # 3. Build the PPTX
    width_cm, height_cm = _PAGE_SIZE_CM.get(page_format, _PAGE_SIZE_CM["presentation_16_9"])
    slide_width_emu = _cm_to_emu(width_cm)
    slide_height_emu = _cm_to_emu(height_cm)

    prs = Presentation()
    prs.slide_width = Emu(slide_width_emu)
    prs.slide_height = Emu(slide_height_emu)

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    for img in slide_images:
        slide = prs.slides.add_slide(blank_layout)

        # Save image to an in-memory buffer
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        buf.seek(0)

        # Fill entire slide with the image
        slide.shapes.add_picture(
            buf,
            left=Emu(0),
            top=Emu(0),
            width=Emu(slide_width_emu),
            height=Emu(slide_height_emu),
        )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
